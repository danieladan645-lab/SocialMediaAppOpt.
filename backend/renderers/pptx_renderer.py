from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from models import AuditData

# === COLORS ===
CORAL       = (232, 96,  76)
TEAL        = (61,  191, 191)
WARM_WHITE  = (255, 248, 245)
DARK_TEXT   = (44,  44,  44)
WHITE       = (255, 255, 255)
MUTED_GRAY  = (217, 217, 217)
GOLD        = (245, 166, 35)
LIGHT_TEAL  = (220, 245, 245)
LIGHT_CORAL = (252, 230, 226)

COLOR_MAP = {
    "coral": CORAL,
    "teal":  TEAL,
    "gold":  GOLD,
    "gray":  MUTED_GRAY,
}


# === HELPERS ===

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(*color)


def _textbox(slide, left, top, width, height, text, size, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False, bullet_char=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = (bullet_char + "  " + text) if bullet_char else text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(*color)
    return tf


def _add_para(tf, text, size, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
              space_before=4, bullet_char=None):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = (bullet_char + "  " + text) if bullet_char else text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(*color)
    return p


def _rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(*fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = _rgb(*line_color)
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(*line_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _score_badge(slide, score_text, left, top, fill=GOLD, text_color=WHITE):
    shape = _rounded_rect(slide, left, top, 1.6, 0.65, fill)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = score_text
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = _rgb(*text_color)


def _color_bar(slide, left, top, height, fill_color):
    bar = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.08), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*fill_color)
    bar.line.fill.background()


def _table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height)).table


def _fmt_cell(cell, text, size=11, bold=False, fill=None, text_color=DARK_TEXT,
              align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(*text_color)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*fill)


def _slide_title(slide, text, bg_color=None):
    if bg_color:
        _rect(slide, 0, 0, 13.33, 0.72, fill_color=bg_color)
    _textbox(slide, 0.3, 0.1, 12.5, 0.55, text, 22, bold=True,
             color=WHITE if bg_color else CORAL)


# === SLIDE BUILDERS ===

def _slide_01_title(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, CORAL)
    _textbox(slide, 1.0, 1.8, 11.0, 1.2,
             f"{d.handle} Instagram Profile Audit", 38, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 3.1, 11.0, 0.5,
             "A Professional Content & Conversion Audit", 18,
             color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    _textbox(slide, 1.0, 3.7, 11.0, 0.4,
             f"Prepared: {d.date}  |  {d.display_name}", 13,
             color=WARM_WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 0, 6.9, 13.33, 0.6, fill_color=TEAL)
    snap = d.profile_snapshot
    _textbox(slide, 0, 6.95, 13.33, 0.4,
             f"{snap.followers} Followers  •  {snap.platforms}", 11,
             color=WHITE, align=PP_ALIGN.CENTER)


def _slide_02_toc(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Table of Contents", bg_color=TEAL)
    left_items = [
        "1.  Executive Summary",
        "2.  Profile Snapshot",
        "3.  Bio Analysis",
        "4.  Visual Grid & Aesthetic",
        "5.  Highlights Analysis",
        "6.  Content Strategy",
    ]
    right_start = 7
    right_items = []
    for i, p in enumerate(d.personas):
        right_items.append(f"{right_start + i}.  Customer Persona {i+1} — {p.title.split('—')[0].strip()}")
    offset = right_start + len(d.personas)
    right_items += [
        f"{offset}.  Strengths",
        f"{offset+1}.  Weaknesses & Opportunities",
        f"{offset+2}.  Recommendations & Action Plan",
        f"{offset+3}.  Score Summary",
    ]
    tf_l = _textbox(slide, 0.5, 1.0, 6.2, 5.8, left_items[0], 13, color=DARK_TEXT)
    for item in left_items[1:]:
        _add_para(tf_l, item, 13, color=DARK_TEXT, space_before=8)
    tf_r = _textbox(slide, 7.0, 1.0, 6.2, 5.8, right_items[0], 13, color=DARK_TEXT)
    for item in right_items[1:]:
        _add_para(tf_r, item, 13, color=DARK_TEXT, space_before=8)


def _slide_03_executive(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Executive Summary", bg_color=CORAL)
    _score_badge(slide, d.overall_score, 5.85, 1.0, fill=GOLD, text_color=WHITE)
    _textbox(slide, 4.5, 1.75, 4.3, 0.3, "Overall Audit Score", 11,
             color=DARK_TEXT, align=PP_ALIGN.CENTER, italic=True)
    _textbox(slide, 0.5, 2.3, 3.5, 0.4, "Quick Wins", 14, bold=True, color=TEAL)
    tf = _textbox(slide, 0.5, 2.8, 12.3, 3.5, "", 1, color=DARK_TEXT)
    tf.paragraphs[0].runs
    for qw in d.quick_wins:
        _add_para(tf, qw, 13, color=DARK_TEXT, bullet_char="✔", space_before=10)
    _textbox(slide, 0.5, 6.6, 12.3, 0.4,
             "Full breakdown follows on subsequent slides.", 10,
             color=MUTED_GRAY, italic=True)


def _slide_04_profile_snapshot(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Profile Snapshot", bg_color=CORAL)
    snap = d.profile_snapshot
    tbl = _table(slide, 6, 2, 0.5, 1.0, 8.5, 3.2)
    for i, h in enumerate(["Attribute", "Detail"]):
        _fmt_cell(tbl.cell(0, i), h, size=12, bold=True,
                  fill=TEAL, text_color=WHITE, align=PP_ALIGN.CENTER)
    rows_data = [
        ("Username", snap.username),
        ("Display Name", snap.display_name),
        ("Followers", snap.followers),
        ("Bio Link", snap.bio_link),
        ("Profile Photo", snap.photo_desc),
    ]
    for r, (label, val) in enumerate(rows_data, start=1):
        fill = LIGHT_TEAL if r % 2 == 0 else WHITE
        _fmt_cell(tbl.cell(r, 0), label, size=11, bold=True, fill=fill, text_color=DARK_TEXT)
        _fmt_cell(tbl.cell(r, 1), val, size=11, fill=fill, text_color=DARK_TEXT)

    card_colors = [CORAL, TEAL, GOLD]
    for i, card in enumerate(d.stat_cards[:3]):
        top = 1.0 + i * 1.3
        _rounded_rect(slide, 9.5, top, 3.3, 1.1, card_colors[i])
        _textbox(slide, 9.5, top + 0.05, 3.3, 0.4, card.value, 28, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(slide, 9.5, top + 0.5, 3.3, 0.4, card.label, 11,
                 color=WHITE, align=PP_ALIGN.CENTER)


def _slide_05_bio_analysis(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Bio Analysis", bg_color=CORAL)
    for i, block in enumerate(d.bio_blocks[:3]):
        color = COLOR_MAP.get(block.color_key, CORAL)
        top = 1.1 + i * 1.8
        _color_bar(slide, 0.3, top, 1.5, color)
        _textbox(slide, 0.55, top, 2.5, 0.4, block.label, 13, bold=True, color=color)
        _score_badge(slide, block.score_str, 3.2, top, fill=color, text_color=WHITE)
        _textbox(slide, 0.55, top + 0.45, 12.5, 1.1, block.body, 11, color=DARK_TEXT)


def _slide_06_visual_grid(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Visual Grid & Aesthetic", bg_color=CORAL)
    if d.grid_bullets:
        tf = _textbox(slide, 0.5, 1.0, 7.8, 5.5, "•  " + d.grid_bullets[0], 12, color=DARK_TEXT)
        for b in d.grid_bullets[1:]:
            _add_para(tf, b, 12, bullet_char="•", color=DARK_TEXT, space_before=12)
    _textbox(slide, 9.0, 1.0, 4.0, 0.35, "Content Grid Sample", 11,
             bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    for idx in range(6):
        col = idx % 3
        row = idx // 3
        left = 9.0 + col * 1.33
        top = 1.4 + row * 2.0
        _rect(slide, left, top, 1.25, 1.85, fill_color=MUTED_GRAY, line_color=WHITE)
        _textbox(slide, left, top + 0.7, 1.25, 0.5,
                 f"Post {idx+1}", 9, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    _score_badge(slide, d.grid_score, 5.5, 6.1, fill=CORAL, text_color=WHITE)


def _slide_07_highlights(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Highlights Analysis", bg_color=CORAL)
    rows = d.highlight_rows[:5]
    tbl = _table(slide, len(rows) + 1, 3, 0.5, 1.0, 12.3, 4.2)
    for i, h in enumerate(["Highlight Name", "Cover Quality", "Recommendation"]):
        _fmt_cell(tbl.cell(0, i), h, size=12, bold=True,
                  fill=TEAL, text_color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        fill = LIGHT_TEAL if r % 2 == 0 else WHITE
        _fmt_cell(tbl.cell(r, 0), row.name, size=11, fill=fill, text_color=DARK_TEXT)
        _fmt_cell(tbl.cell(r, 1), row.cover_quality, size=11, fill=fill, text_color=DARK_TEXT)
        _fmt_cell(tbl.cell(r, 2), row.recommendation, size=11, fill=fill, text_color=DARK_TEXT)
    _textbox(slide, 0.5, 5.4, 12.3, 0.5, d.highlight_action, 11, color=DARK_TEXT, italic=True)
    _score_badge(slide, d.highlight_score, 5.85, 5.9, fill=CORAL, text_color=WHITE)


def _slide_08_content_strategy(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Content Strategy", bg_color=CORAL)
    chart_left, chart_top, max_width = 0.5, 1.2, 6.0
    for i, bar in enumerate(d.content_bars[:3]):
        color = COLOR_MAP.get(bar.color_key, CORAL)
        bar_top = chart_top + i * 0.85
        bar_w = max_width * bar.pct / 100
        _rect(slide, chart_left, bar_top, bar_w, 0.55, fill_color=color)
        _textbox(slide, chart_left + bar_w + 0.15, bar_top + 0.05,
                 1.5, 0.45, f"{bar.pct}%", 13, bold=True, color=color)
        _textbox(slide, chart_left, bar_top - 0.32, 5.5, 0.32,
                 bar.label, 10, color=DARK_TEXT, bold=True)
    if d.content_bullets:
        tf = _textbox(slide, 7.0, 1.2, 6.0, 5.0, "•  " + d.content_bullets[0], 12, color=DARK_TEXT)
        for b in d.content_bullets[1:]:
            _add_para(tf, b, 12, bullet_char="•", color=DARK_TEXT, space_before=12)
    _score_badge(slide, d.content_score, 5.85, 6.0, fill=CORAL, text_color=WHITE)


def _persona_slide(prs, persona, slide_num: int):
    border_color = COLOR_MAP.get(persona.color_key, CORAL)
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, f"Customer Persona {slide_num}: {persona.title}", bg_color=border_color)
    _rounded_rect(slide, 0.4, 1.0, 12.5, 5.8, WARM_WHITE, line_color=border_color)
    fields = [
        ("Name & Age", persona.name),
        ("Location", persona.location),
        ("Goals", persona.goals),
        ("Pain Points", persona.pain_points),
        ("How They Found Profile", persona.found_via),
        ("First Impression", persona.first_impression),
        ("What They Need", persona.needs),
        ("Conversion Path", persona.conversion),
    ]
    col_w_label, col_w_val = 3.5, 8.5
    left_label, left_val, top_start, row_h = 0.55, 4.2, 1.15, 0.62
    for i, (label, val) in enumerate(fields):
        top = top_start + i * row_h
        fill = LIGHT_TEAL if i % 2 == 0 else WHITE
        _rect(slide, left_label - 0.1, top, col_w_label + 0.1, row_h - 0.05, fill_color=fill)
        _rect(slide, left_val - 0.1, top, col_w_val + 0.3, row_h - 0.05, fill_color=fill)
        _textbox(slide, left_label, top + 0.05, col_w_label, row_h - 0.1,
                 label, 10, bold=True, color=border_color)
        _textbox(slide, left_val, top + 0.05, col_w_val, row_h - 0.1,
                 val, 10, color=DARK_TEXT)


def _slide_strengths(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Strengths", bg_color=TEAL)
    for i, s in enumerate(d.strengths[:5]):
        top = 1.1 + i * 1.1
        _rounded_rect(slide, 0.4, top, 0.5, 0.5, TEAL)
        _textbox(slide, 0.48, top + 0.08, 0.35, 0.35, "✓", 16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(slide, 1.1, top, 4.0, 0.4, s.title, 13, bold=True, color=TEAL)
        _textbox(slide, 1.1, top + 0.42, 11.8, 0.55, s.body, 11, color=DARK_TEXT)


def _slide_weaknesses(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Weaknesses & Opportunities", bg_color=CORAL)
    _rect(slide, 0.4, 0.9, 6.0, 0.4, fill_color=LIGHT_CORAL)
    _textbox(slide, 0.5, 0.92, 5.8, 0.36, "Weaknesses", 13, bold=True, color=CORAL)
    if d.weaknesses:
        tf_w = _textbox(slide, 0.5, 1.4, 5.8, 5.5, d.weaknesses[0], 12,
                        bullet_char="●", color=CORAL)
        for w in d.weaknesses[1:]:
            _add_para(tf_w, w, 12, bullet_char="●", color=CORAL, space_before=14)
    _rect(slide, 7.0, 0.9, 6.0, 0.4, fill_color=LIGHT_TEAL)
    _textbox(slide, 7.1, 0.92, 5.8, 0.36, "Opportunities", 13, bold=True, color=TEAL)
    if d.opportunities:
        tf_o = _textbox(slide, 7.1, 1.4, 5.8, 5.5, d.opportunities[0], 12,
                        bullet_char="●", color=TEAL)
        for o in d.opportunities[1:]:
            _add_para(tf_o, o, 12, bullet_char="●", color=TEAL, space_before=14)


def _slide_recommendations(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Recommendations & Action Plan", bg_color=CORAL)
    recs = d.recommendations
    tbl = _table(slide, len(recs) + 1, 3, 0.4, 1.0, 12.5, 5.8)
    for i, h in enumerate(["Priority", "Action", "Expected Impact"]):
        _fmt_cell(tbl.cell(0, i), h, size=12, bold=True,
                  fill=DARK_TEXT, text_color=WHITE, align=PP_ALIGN.CENTER)
    p_color_map = {"HIGH": CORAL, "MED": TEAL, "LOW": MUTED_GRAY}
    p_fill_map  = {"HIGH": LIGHT_CORAL, "MED": LIGHT_TEAL, "LOW": WHITE}
    for r, rec in enumerate(recs, start=1):
        p_color = p_color_map.get(rec.priority, MUTED_GRAY)
        fill = p_fill_map.get(rec.priority, WHITE)
        _fmt_cell(tbl.cell(r, 0), rec.priority, size=11, bold=True,
                  fill=fill, text_color=p_color, align=PP_ALIGN.CENTER)
        _fmt_cell(tbl.cell(r, 1), rec.action, size=10, fill=fill, text_color=DARK_TEXT)
        _fmt_cell(tbl.cell(r, 2), rec.impact, size=10, fill=fill, text_color=DARK_TEXT)


def _slide_scorecard(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, WARM_WHITE)
    _slide_title(slide, "Score Summary", bg_color=TEAL)
    rows = d.score_rows
    tbl = _table(slide, len(rows) + 1, 2, 2.5, 1.1, 8.5, 5.2)
    _fmt_cell(tbl.cell(0, 0), "Category", size=12, bold=True,
              fill=TEAL, text_color=WHITE, align=PP_ALIGN.CENTER)
    _fmt_cell(tbl.cell(0, 1), "Score", size=12, bold=True,
              fill=TEAL, text_color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        is_overall = row.category.upper() == "OVERALL"
        fill = GOLD if is_overall else (LIGHT_TEAL if r % 2 == 0 else LIGHT_CORAL)
        text_color = WHITE if is_overall else DARK_TEXT
        size = 14 if is_overall else 12
        _fmt_cell(tbl.cell(r, 0), row.category, size=size, bold=is_overall,
                  fill=fill, text_color=text_color)
        _fmt_cell(tbl.cell(r, 1), row.score_str, size=size, bold=is_overall,
                  fill=fill, text_color=text_color, align=PP_ALIGN.CENTER)
    _textbox(slide, 2.5, 6.4, 8.5, 0.4,
             f"Audit conducted: {d.date}  |  {d.handle} Instagram Profile", 10,
             color=MUTED_GRAY, align=PP_ALIGN.CENTER, italic=True)


def _slide_closing(prs, d: AuditData):
    slide = _add_slide(prs)
    _set_bg(slide, TEAL)
    _textbox(slide, 1.0, 2.2, 11.3, 0.9, "Thank You", 40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 3.2, 11.3, 0.55,
             f"{d.handle} — Keep building.", 18,
             color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    _textbox(slide, 1.0, 3.9, 11.3, 0.4,
             f"Audit prepared for: {d.display_name}  |  {d.date}", 12,
             color=WARM_WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 0, 6.9, 13.33, 0.6, fill_color=CORAL)
    _textbox(slide, 0, 6.95, 13.33, 0.4,
             d.profile_snapshot.platforms, 11, color=WHITE, align=PP_ALIGN.CENTER)


# === ENTRY POINT ===

def build_pptx(d: AuditData) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_01_title(prs, d)
    _slide_02_toc(prs, d)
    _slide_03_executive(prs, d)
    _slide_04_profile_snapshot(prs, d)
    _slide_05_bio_analysis(prs, d)
    _slide_06_visual_grid(prs, d)
    _slide_07_highlights(prs, d)
    _slide_08_content_strategy(prs, d)
    for i, persona in enumerate(d.personas, start=1):
        _persona_slide(prs, persona, i)
    _slide_strengths(prs, d)
    _slide_weaknesses(prs, d)
    _slide_recommendations(prs, d)
    _slide_scorecard(prs, d)
    _slide_closing(prs, d)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
